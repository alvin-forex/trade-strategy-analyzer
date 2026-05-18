#!/usr/bin/env python3
"""
Trade Strategy Analyzer — 系統說明書 PPT 生成器
包含：每頁功能、使用方法、流程範例、計算原理、風險分析
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ============================================================
# 配色方案
# ============================================================
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 深藍黑背景
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)        # 卡片背景
ACCENT_RED = RGBColor(0xE9, 0x45, 0x60)     # 主色紅
ACCENT_BLUE = RGBColor(0x0F, 0x3D, 0x96)    # 輔助藍
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)     # 白色文字
TEXT_GRAY = RGBColor(0xA0, 0xA8, 0xB8)      # 灰色文字
GREEN = RGBColor(0x4C, 0xAF, 0x50)          # 綠色
YELLOW = RGBColor(0xFF, 0xC1, 0x07)         # 黃色
ORANGE = RGBColor(0xFF, 0x98, 0x00)         # 橙色
RED_DANGER = RGBColor(0xFF, 0x57, 0x22)     # 危險紅

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# Helper Functions
# ============================================================
def add_bg(slide, color=BG_DARK):
    """Add solid background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=13, default_color=TEXT_WHITE):
    """Add multiple formatted lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, size, color, bold = line_info, default_size, default_color, False
        else:
            text = line_info.get('text', '')
            size = line_info.get('size', default_size)
            color = line_info.get('color', default_color)
            bold = line_info.get('bold', False)
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Segoe UI'
        p.space_after = Pt(4)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_arrow_box(slide, left, top, width, height, text, fill_color=ACCENT_BLUE, text_color=TEXT_WHITE, font_size=11):
    """Add a rounded rect with centered text (for flow diagrams)."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = 'Segoe UI'
    return shape

def add_section_header(slide, title, subtitle=""):
    """Add consistent section header."""
    # Top accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_RED
    shape.line.fill.background()
    
    add_text_box(slide, 0.6, 0.3, 12, 0.6, title, font_size=28, color=ACCENT_RED, bold=True)
    if subtitle:
        add_text_box(slide, 0.6, 0.85, 12, 0.4, subtitle, font_size=14, color=TEXT_GRAY)

def add_risk_card(slide, left, top, width, height, title, risks):
    """Add a risk analysis card."""
    add_rounded_rect(slide, left, top, width, height, RGBColor(0x2A, 0x15, 0x15))
    add_text_box(slide, left + 0.15, top + 0.08, width - 0.3, 0.35, f"⚠️ {title}", font_size=13, color=RED_DANGER, bold=True)
    lines = [{'text': f"• {r}", 'size': 11, 'color': RGBColor(0xFF, 0xCC, 0xCC)} for r in risks]
    add_multi_text(slide, left + 0.15, top + 0.4, width - 0.3, height - 0.5, lines)

# ============================================================
# SLIDE 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Big accent bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.8))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, 1, 1.2, 11, 0.5, "TRADE STRATEGY ANALYZER", font_size=16, color=TEXT_GRAY, bold=True)
add_text_box(slide, 1, 2.7, 11, 1.2, "系統架構說明書", font_size=48, color=TEXT_WHITE, bold=True)
add_text_box(slide, 1, 3.9, 11, 0.6, "每頁功能 · 使用方法 · 計算原理 · 流程演示 · 風險分析", font_size=20, color=RGBColor(0xCC, 0xDD, 0xFF))

add_text_box(slide, 1, 5.8, 5, 0.4, f"版本 v0.7  |  {datetime.date.today().strftime('%Y-%m-%d')}", font_size=14, color=TEXT_GRAY)
add_text_box(slide, 1, 6.2, 5, 0.4, "丁蟹 🦀 + Alvin", font_size=14, color=TEXT_GRAY)

# ============================================================
# SLIDE 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📖 目錄 — 20 個功能頁面一覽")

tabs = [
    ("📈", "收益曲線", "Equity Curve"), ("🎯", "質量評分", "Quality Scoring"),
    ("💱", "貨幣對", "Symbol Analysis"), ("📚", "層數分析", "Layer Analysis"),
    ("🕐", "時段分析", "Session Analysis"), ("⏱️", "持倉時間", "Holding Time"),
    ("↔️", "方向分析", "Direction Analysis"), ("🌐", "市場語境", "Market Context"),
    ("🔄", "市況×策略", "Market Regime"), ("⚖️", "策略對比", "Strategy Compare"),
    ("🔍", "貨幣穿透", "Symbol Drill"), ("🏆", "Copy 推薦", "Copy Matrix"),
    ("📏", "波幅分析", "Volatility"), ("🎲", "馬丁風險", "Martin Risk"),
    ("📐", "TP/SL", "TP/SL Analysis"), ("💰", "Copy 模擬", "Copy Simulation"),
    ("📋", "倉位明細", "Positions"), ("🔬", "馬丁剖析V3", "Martin Autopsy"),
    ("⚙️", "SET 參數", "SET Params"), ("📁", "歷史記錄", "Archive"),
]

for i, (icon, cn, en) in enumerate(tabs):
    col = i % 4
    row = i // 4
    x = 0.6 + col * 3.1
    y = 1.5 + row * 1.15
    add_rounded_rect(slide, x, y, 2.9, 0.95, BG_CARD)
    add_text_box(slide, x + 0.15, y + 0.08, 2.6, 0.35, f"{icon} {cn}", font_size=15, color=ACCENT_RED, bold=True)
    add_text_box(slide, x + 0.15, y + 0.48, 2.6, 0.3, en, font_size=11, color=TEXT_GRAY)

# ============================================================
# SLIDE 3: System Architecture Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "🏗️ 系統架構總覽", "數據流 + 前後端架構")

# Data flow diagram
flow_steps = [
    ("CSV 上傳\nparseCSV()", ACCENT_BLUE),
    ("SET 上傳\nparseSET()", ACCENT_BLUE),
    ("層級偵測\nbuildLayerMapping()", RGBColor(0x1B, 0x5E, 0x20)),
    ("倉位重構\nbuildPositions()", RGBColor(0x1B, 0x5E, 0x20)),
    ("多維計算\ncalcStats/Score...", ORANGE),
    ("20 個 Tab\n渲染展示", ACCENT_RED),
]

for i, (text, color) in enumerate(flow_steps):
    x = 0.5 + i * 2.1
    add_arrow_box(slide, x, 1.6, 1.9, 0.8, text, color, font_size=10)
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + 1.9, 1.85, 0.25, 0.3, "→", font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Architecture cards
add_rounded_rect(slide, 0.5, 3.0, 4.0, 4.0, BG_CARD)
add_multi_text(slide, 0.7, 3.1, 3.6, 3.8, [
    {'text': '🖥️ 前端（單文件 HTML）', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '• 部署：GitHub Pages', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 框架：純 JavaScript（零依賴）', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 圖表：SVG 折線圖', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 數據：用戶上傳 CSV + SET', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 計算：全部客戶端本地完成', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 離線可用', 'size': 12, 'color': GREEN},
])

add_rounded_rect(slide, 4.8, 3.0, 4.0, 4.0, BG_CARD)
add_multi_text(slide, 5.0, 3.1, 3.6, 3.8, [
    {'text': '⚙️ 後端（Python）', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '• DDE v3 評分引擎', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• Lot-Based 層級偵測', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 69 Signals 批量分析', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• HTML 報告自動生成', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• FastAPI (localhost:8787)', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• AlgoForest Scraper', 'size': 12, 'color': TEXT_WHITE},
])

add_rounded_rect(slide, 9.1, 3.0, 3.8, 4.0, BG_CARD)
add_multi_text(slide, 9.3, 3.1, 3.4, 3.8, [
    {'text': '💾 數據存儲', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '• SQLite: analysis_history.db', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• JSON: batch/signal configs', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• CSV: 69 個信號交易數據', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• SET: EA 策略設定檔', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• localStorage: 前端存檔', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• MT4 .hst: D1 市場數據', 'size': 12, 'color': TEXT_WHITE},
])

# ============================================================
# SLIDE 4: Input Data Specs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📥 輸入數據規格", "CSV 交易數據 + SET 策略設定檔")

# CSV columns
add_rounded_rect(slide, 0.5, 1.4, 6.2, 5.8, BG_CARD)
add_text_box(slide, 0.7, 1.5, 5.8, 0.4, "📄 CSV 交易數據（AlgoForest 下載）", font_size=15, color=ACCENT_RED, bold=True)

csv_fields = [
    "Open Time / Close Time — 開平倉時間 → 時段分析、持倉時間",
    "Type — buy/sell/balance/credit → 只處理 buy/sell",
    "Lots — 手數 → Lot-Based 層級偵測核心",
    "Symbol — 貨幣對 → 分組分析",
    "Open Price / Close Price — 開平倉價 → 計算實際點差",
    "Net Pips — 淨賺點數 → 方向判斷、盈虧",
    "Net Profit — 淨盈虧 → 真實收益計算",
    "Max Pips — 最大浮盈 → 入市時機、TP 建議",
    "Max Loss Pips — 最大浮虧 → 風險承受、SL 建議",
    "Commission / Swap — 手續費/隔夜息 → 成本分析",
    "Magic Number — 88=BUY, 77=SELL → 策略識別",
    "Holding Time — 持倉時間 → 效率分析",
]
lines = [{'text': f"  {f}", 'size': 11, 'color': TEXT_WHITE} for f in csv_fields]
add_multi_text(slide, 0.7, 2.0, 5.8, 5.0, lines)

# SET info
add_rounded_rect(slide, 7.0, 1.4, 5.8, 5.8, BG_CARD)
add_text_box(slide, 7.2, 1.5, 5.4, 0.4, "⚙️ SET 策略設定檔", font_size=15, color=ACCENT_RED, bold=True)

set_info = [
    {'text': 'EA 家族對照表：', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '  Dragon Wave (DW) — LotMul×2.5 倍投，8 層', 'size': 11, 'color': RGBColor(0xCE, 0x93, 0xD8)},
    {'text': '  SMA — lotExp 指數遞增 + pipstep，7-15 層', 'size': 11, 'color': RGBColor(0xA5, 0xD6, 0xA7)},
    {'text': '  MKD — PipStep 網格，6-10 層', 'size': 11, 'color': RGBColor(0xFF, 0xCC, 0x80)},
    {'text': '  S10 — 固定 lotSize 平注碼，MaxBuyCount=10', 'size': 11, 'color': RGBColor(0x90, 0xCA, 0xF9)},
    {'text': '  Flash — CheckLevels，11 層', 'size': 11, 'color': RGBColor(0xF4, 0x8F, 0xB1)},
    {'text': '  GEM — 無馬丁', 'size': 11, 'color': RGBColor(0xB0, 0xBE, 0xC5)},
    {'text': '', 'size': 8, 'color': TEXT_WHITE},
    {'text': '層級偵測方法（v0.7 Lot-Based）：', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '  1️⃣ Primary：signal_lot_mapping.json', 'size': 11, 'color': TEXT_WHITE},
    {'text': '      → SET 檔 lot→level 對照表直接映射', 'size': 11, 'color': TEXT_GRAY},
    {'text': '  2️⃣ Fallback：CSV Lots 唯一值推算', 'size': 11, 'color': TEXT_WHITE},
    {'text': '      → 無 SET 時按 unique lot 排序', 'size': 11, 'color': TEXT_GRAY},
    {'text': '  3️⃣ AutoLot 標記：lots >> SET layers', 'size': 11, 'color': TEXT_WHITE},
    {'text': '', 'size': 8, 'color': TEXT_WHITE},
    {'text': '⚠️ 過濾規則：Type=balance/credit 必須排除！', 'size': 12, 'color': RED_DANGER, 'bold': True},
]
add_multi_text(slide, 7.2, 2.0, 5.4, 5.0, set_info)

# ============================================================
# SLIDE 5-6: DDE v3 Scoring System
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "🧮 DDE v3 評分系統 — 計算原理", "Copy Trade 場景專用：評估「跟單」可行性同回報")

# Formula breakdown
add_rounded_rect(slide, 0.5, 1.4, 12.3, 1.4, ACCENT_BLUE)
add_text_box(slide, 0.7, 1.5, 11.9, 0.5, "DDE v3 Score = Trigger Rate (40%) + Alpha Capture Profit (40%) + DDE 回撤效率 (20%)", 
             font_size=20, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.1, 11.9, 0.5, "評級：⭐⭐⭐⭐ ≥80  |  ⭐⭐⭐ ≥60  |  ⭐⭐ ≥40  |  ⭐ <40", 
             font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Three pillars
for i, (title, weight, formula, example) in enumerate([
    ("指標 1：Trigger Rate\n觸發率", "40%",
     "CoP: 觸發數/盈利交易數\ntrigger = Max Pips ≥ wait_pips\nwait: 5, 10, 15, 20 pips\n\nCoL: 觸發數/虧損交易數\ntrigger = |Max Loss Pips| ≥ wait_pips\nwait: 10, 15, 20, 25 pips\n\nScore = min(rate×100, 100) × 0.4",
     "例：Signal 22200 AUDCAD L1\n盈利交易 100 筆\nWait=5 時 95 筆觸發\n→ rate = 95%\n→ Score = 95 × 0.4 = 38.0"),
    ("指標 2：Alpha Capture\n利潤捕捉分", "40%",
     "動態百分位評分 (0-120)\n\nbaseline = max(P50, global_P25)\nfloor = $5.00（垃圾過濾）\n\n≥ P75: 100 + bonus (cap 120)\n≥ P50: 70 + 線性插值(70-100)\n≥ baseline: 線性插值(0-70)\n< baseline: 0",
     "例：avg_profit = $45\nP50 = $30, P75 = $60\n→ 介於 P50-P75 之間\n→ 70 + (45-30)/(60-30)×30\n   = 70 + 15 = 85\n→ Score = 85 × 0.4 = 34.0"),
    ("指標 3：DDE\n回撤效率", "20%",
     "dd_ratio = |max_loss_pips|/profit_pips\n每筆交易計算，cap at 2.0\navg_dd_ratio = mean(所有 dd_ratio)\n\nScore = max(0, 100 - 50 × avg) × 0.2\n\n取代舊 ETE（分辨力≈0）\nDDE 分佈：0-100，中位數 76",
     "例：5 筆交易\nratios: 0.3, 0.5, 0.8, 1.2, 0.4\navg = 0.64\n→ 100 - 50×0.64 = 68\n→ Score = 68 × 0.2 = 13.6"),
]):
    x = 0.5 + i * 4.15
    add_rounded_rect(slide, x, 3.1, 3.95, 4.2, BG_CARD)
    add_text_box(slide, x + 0.15, 3.15, 3.65, 0.5, title, font_size=13, color=ACCENT_RED, bold=True)
    add_text_box(slide, x + 0.15, 3.65, 1.8, 0.3, f"權重：{weight}", font_size=11, color=YELLOW, bold=True)
    add_text_box(slide, x + 0.15, 3.95, 3.65, 2.2, formula, font_size=10, color=TEXT_WHITE)
    add_rounded_rect(slide, x + 0.15, 6.2, 3.65, 0.95, RGBColor(0x1B, 0x2A, 0x1B))
    add_text_box(slide, x + 0.2, 6.25, 3.55, 0.85, example, font_size=9, color=GREEN)

# ============================================================
# SLIDE 6: DDE v3 Flow Example
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "🔄 DDE v3 完整分析流程 — 以 Signal 22200 為例")

# Flow chart
steps = [
    ("📥 下載 CSV\nalgoforest_scraper.py", ACCENT_BLUE),
    ("🔍 Lot 層級偵測\nassign_lot_level()", RGBColor(0x1B, 0x5E, 0x20)),
    ("📊 按 CCY×Level\n分組計算", ORANGE),
    ("🎯 CoP 評分\nWait 5/10/15/20", ACCENT_RED),
    ("🎯 CoL 評分\nWait 10/15/20/25", ACCENT_RED),
    ("📋 加總 Avg Score\n生成排名", RGBColor(0x4A, 0x14, 0x8C)),
]

for i, (text, color) in enumerate(steps):
    x = 0.3 + i * 2.15
    add_arrow_box(slide, x, 1.4, 2.0, 0.8, text, color, font_size=10)
    if i < len(steps) - 1:
        add_text_box(slide, x + 2.0, 1.65, 0.2, 0.3, "→", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Example calculation
add_rounded_rect(slide, 0.5, 2.5, 12.3, 4.7, BG_CARD)
add_text_box(slide, 0.7, 2.6, 12, 0.4, "📝 計算範例：Signal 22200 / AUDCAD / L1 / CoP Wait=10", font_size=14, color=ACCENT_RED, bold=True)

example_lines = [
    {'text': 'Step 1：篩選數據', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '  • AUDCAD L1 交易共 45 筆，其中盈利 35 筆（net_profit > 0）', 'size': 11, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Step 2：Trigger Rate（權重 40%）', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '  • CoP Wait=10：檢查每筆盈利交易是否 Max Pips ≥ 10', 'size': 11, 'color': TEXT_WHITE},
    {'text': '  • 結果：32/35 筆觸發 → trigger_rate = 91.4%', 'size': 11, 'color': TEXT_WHITE},
    {'text': '  • Score = min(91.4, 100) × 0.4 = 36.6', 'size': 11, 'color': GREEN},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Step 3：Alpha Capture（權重 40%）', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '  • 觸發交易平均盈利：$52.30', 'size': 11, 'color': TEXT_WHITE},
    {'text': '  • Signal P50=$35, P75=$65 → 介於兩者之間', 'size': 11, 'color': TEXT_WHITE},
    {'text': '  • Score = (70 + (52.3-35)/(65-35)×30) × 0.4 = 70 + 17.3 = 87.3 × 0.4 = 34.9', 'size': 11, 'color': GREEN},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Step 4：DDE（權重 20%）', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '  • 每筆：dd_ratio = |Max Loss Pips| / Profit Pips，cap 2.0', 'size': 11, 'color': TEXT_WHITE},
    {'text': '  • 平均 dd_ratio = 0.55 → Score = (100 - 50×0.55) × 0.2 = 72.5 × 0.2 = 14.5', 'size': 11, 'color': GREEN},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Step 5：總分 = 36.6 + 34.9 + 14.5 = 86.0 → ⭐⭐⭐⭐ 建議 Copy ✅', 'size': 14, 'color': ACCENT_RED, 'bold': True},
]
add_multi_text(slide, 0.7, 3.0, 11.9, 4.0, example_lines)

# ============================================================
# SLIDE 7: Tab 1-4 — 收益曲線 / 質量評分 / 貨幣對 / 層數分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📊 頁面功能 1-4：基礎分析 Tabs", "收益曲線 → 質量評分 → 貨幣對 → 層數分析")

# Tab 1: 收益曲線
add_rounded_rect(slide, 0.4, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 0.55, 1.35, 2.8, 0.35, "📈 收益曲線 (Equity Curve)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 0.55, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'SVG 折線圖展示累積盈虧', 'size': 10, 'color': TEXT_WHITE},
    {'text': '正/負區域著色（綠盈利紅虧損）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '關鍵統計卡片（總盈利、勝率等）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '使用方法：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '1. 上傳 CSV 自動渲染', 'size': 10, 'color': TEXT_WHITE},
    {'text': '2. 第一眼睇總體盈虧趨勢', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '計算：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '累積 net_profit 逐筆相加', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'drawLine() → SVG polyline', 'size': 10, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 曲線上升可能因馬丁加倉', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 掩蓋深層風險（只睇結果）', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 需配合層數分析使用', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 2: 質量評分
add_rounded_rect(slide, 3.6, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 3.75, 1.35, 2.8, 0.35, "🎯 質量評分 (Quality Score)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 3.75, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'Entry Score（入場質素）', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Strategy Score（策略管理）', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Final = Entry×0.4 + Strategy×0.6', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Entry Score 計算：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '方向(35%) — 趨勢對唔對', 'size': 10, 'color': TEXT_WHITE},
    {'text': '時機(35%) — 浮盈 vs 浮虧', 'size': 10, 'color': TEXT_WHITE},
    {'text': '初始回撤(30%) — 捱打幾耐', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'L4+ timing ×1.2 馬丁修正', 'size': 10, 'color': ORANGE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Strategy Score 計算：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '回歸率(30%) — 深層自救能力', 'size': 10, 'color': TEXT_WHITE},
    {'text': '出場效率(25%) — 賺到幾多', 'size': 10, 'color': TEXT_WHITE},
    {'text': '風控(20%)+品質(15%)+成本(10%)', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 評分基於歷史非即時', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 無 Tick Data，精度受限', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 3: 貨幣對
add_rounded_rect(slide, 6.8, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 6.95, 1.35, 2.8, 0.35, "💱 貨幣對 (Symbol Analysis)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 6.95, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '按 PF（Profit Factor）排序', 'size': 10, 'color': TEXT_WHITE},
    {'text': '每個貨幣對：勝率、盈虧、筆數', 'size': 10, 'color': TEXT_WHITE},
    {'text': '點擊展開詳細', 'size': 10, 'color': TEXT_WHITE},
    {'text': '顏色編碼（綠=好、紅=差）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '使用方法：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '1. 上傳 CSV 後自動按 CCY 分組', 'size': 10, 'color': TEXT_WHITE},
    {'text': '2. 按 PF 排序睇邊個最穩定', 'size': 10, 'color': TEXT_WHITE},
    {'text': '3. 展開睇個別交易明細', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '計算：calcSymbolStats()', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'PF = gross_win / gross_loss', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'WR = wins / total × 100%', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• PF 高可能因為少數大贏', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 樣本少時 PF 不穩定', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 4: 層數分析
add_rounded_rect(slide, 10.0, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 10.15, 1.35, 2.8, 0.35, "📚 層數分析 (Layer Analysis)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 10.15, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '馬丁策略核心 — 層數深度分析', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'L1 only → L1-L2 → L1-L3 → L4+', 'size': 10, 'color': TEXT_WHITE},
    {'text': '每層：勝率、盈利、回歸率', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '健康指標：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '✅ L1-only ≥ 50% = 健康', 'size': 10, 'color': GREEN},
    {'text': '⚠️ L4+ ≤ 20% = 可控', 'size': 10, 'color': YELLOW},
    {'text': '✅ L4+ 回歸率 ≥ 70% = 自救力強', 'size': 10, 'color': GREEN},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Lot-Based 偵測（v0.7）：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'buildLayerMapping() → lot→level', 'size': 10, 'color': TEXT_WHITE},
    {'text': '非舊版 pip 硬編碼', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• L4+ 是爆倉高危區', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• AutoLot 信號層級不確定', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 回歸率依賴歷史數據', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# ============================================================
# SLIDE 8: Tab 5-8 — 時段/持倉/方向/市場語境
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📊 頁面功能 5-8：時間與市場分析", "時段分析 → 持倉時間 → 方向分析 → 市場語境")

# Tab 5: 時段分析
add_rounded_rect(slide, 0.4, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 0.55, 1.35, 2.8, 0.35, "🕐 時段分析 (Session)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 0.55, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '按交易時段統計表現差異', 'size': 10, 'color': TEXT_WHITE},
    {'text': '24h 逐小時棒型圖', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '時段定義（HKT）：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '亞洲盤 00-08（波幅細）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '歐洲盤 14-22（波幅大）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '美洲盤 21-05（跨午夜）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '重疊：亞洲優先 0-4', 'size': 10, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '計算：calcSessionStats()', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'calcHourlyStats() → 24格', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 只用 Open Time 分配時段', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 跨時段持倉無法反映', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 6: 持倉時間
add_rounded_rect(slide, 3.6, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 3.75, 1.35, 2.8, 0.35, "⏱️ 持倉時間 (Holding Time)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 3.75, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '5 段持倉時間分組分析', 'size': 10, 'color': TEXT_WHITE},
    {'text': '反映策略類型（日內vs波段）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '分組（用分鐘計）：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '<30m — 剝頭皮型', 'size': 10, 'color': TEXT_WHITE},
    {'text': '30m-1h — 短線', 'size': 10, 'color': TEXT_WHITE},
    {'text': '1-4h — 日內', 'size': 10, 'color': TEXT_WHITE},
    {'text': '4-12h — 波段', 'size': 10, 'color': TEXT_WHITE},
    {'text': '12h+ — 長線', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '計算：calcHoldingStats()', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '用 CSV Holding Time (Hours)', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 長時間持倉 = 高 swap 成本', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 馬丁策略持倉往往偏長', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 7: 方向分析
add_rounded_rect(slide, 6.8, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 6.95, 1.35, 2.8, 0.35, "↔️ 方向分析 (Direction)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 6.95, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '分析 BUY vs SELL 表現差異', 'size': 10, 'color': TEXT_WHITE},
    {'text': '按貨幣對 × 方向 矩陣展示', 'size': 10, 'color': TEXT_WHITE},
    {'text': '識別策略偏多/偏空傾向', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '計算：calcDirectionStats()', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'Magic Number: 88=BUY, 77=SELL', 'size': 10, 'color': TEXT_WHITE},
    {'text': '或用 Type: buy/sell', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '使用場景：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '睇 EA 係咪只識做單邊', 'size': 10, 'color': TEXT_WHITE},
    {'text': '某方向特別弱要留意', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 單方向依賴 = 趨勢反轉高危', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 樣本不均時結果偏頗', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 8: 市場語境
add_rounded_rect(slide, 10.0, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 10.15, 1.35, 2.8, 0.35, "🌐 市場語境 (Market Context)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 10.15, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '結合 MT4 D1 數據顯示同期走勢', 'size': 10, 'color': TEXT_WHITE},
    {'text': '買賣方向標記在價格圖上', 'size': 10, 'color': TEXT_WHITE},
    {'text': '自動判定順勢/逆勢', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '數據來源：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '本地 MT4 .hst → export_hst.py', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'docs/data/ D1 JSON + manifest', 'size': 10, 'color': TEXT_WHITE},
    {'text': '完全離線，無需 API', 'size': 10, 'color': GREEN},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '計算：buildCSVMarketData()', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'renderMarketContext() → SVG', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• D1 精度對 M5/H1 EA 唔夠', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• .hst 數據需手動更新', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 順/逆勢判定為近似', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# ============================================================
# SLIDE 9: Tab 9-12 — 市況×策略 / 策略對比 / 貨幣穿透 / Copy推薦
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📊 頁面功能 9-12：進階分析與 Copy 推薦", "市況×策略 → 策略對比 → 貨幣穿透 → Copy 推薦矩陣")

# Tab 9: 市況×策略
add_rounded_rect(slide, 0.4, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 0.55, 1.35, 2.8, 0.35, "🔄 市況×策略 (Regime×Strategy)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 0.55, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '市況分類 + 策略表現對比', 'size': 10, 'color': TEXT_WHITE},
    {'text': '支援雙 CSV 對比（同 EA 不同 SET）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '市況分類算法：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '1. ATR > 1.5× 均值？ → 高波動🟠', 'size': 10, 'color': ORANGE},
    {'text': '2. |SMA20-50|/mid<0.15% 且 RSI 40-60？', 'size': 10, 'color': TEXT_WHITE},
    {'text': '   → 震盪🟡', 'size': 10, 'color': YELLOW},
    {'text': '3. SMA20>50 且 RSI>50 → 上升🟢', 'size': 10, 'color': GREEN},
    {'text': '4. SMA20<50 且 RSI<50 → 下降🔴', 'size': 10, 'color': RED_DANGER},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 分類基於 SMA/RSI 閾值，主觀', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• D1 分類對短週期 EA 不精確', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 10: 策略對比
add_rounded_rect(slide, 3.6, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 3.75, 1.35, 2.8, 0.35, "⚖️ 策略對比 (Compare)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 3.75, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '上傳第二份 CSV 做並排比較', 'size': 10, 'color': TEXT_WHITE},
    {'text': '按貨幣對、層數、時段逐項對比', 'size': 10, 'color': TEXT_WHITE},
    {'text': '支援市況切片對比', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '使用流程：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '1. 上傳第一份 CSV（基準）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '2. 到此 Tab 上傳第二份', 'size': 10, 'color': TEXT_WHITE},
    {'text': '3. 系統自動並排渲染', 'size': 10, 'color': TEXT_WHITE},
    {'text': '4. 按篩選器縮窄範圍', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 兩份 CSV 時間範圍可能不同', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 市場環境差異未完全消除', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 11: 貨幣穿透
add_rounded_rect(slide, 6.8, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 6.95, 1.35, 2.8, 0.35, "🔍 貨幣穿透 (Symbol Drill)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 6.95, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '選定貨幣對深度分析', 'size': 10, 'color': TEXT_WHITE},
    {'text': '整合 API 查詢歷史分析結果', 'size': 10, 'color': TEXT_WHITE},
    {'text': '可跨 Signal 比較同 CCY', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '使用方法：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '1. 選擇目標貨幣對', 'size': 10, 'color': TEXT_WHITE},
    {'text': '2. renderSymbolDrillResults()', 'size': 10, 'color': TEXT_WHITE},
    {'text': '3. 載入 API 歷史 + 當前數據', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 依賴 localhost:8787 API', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• API 未啟動時功能不可用', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 12: Copy 推薦
add_rounded_rect(slide, 10.0, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 10.15, 1.35, 2.8, 0.35, "🏆 Copy 推薦 (Copy Matrix)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 10.15, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '貨幣對 × 層數 推薦矩陣', 'size': 10, 'color': TEXT_WHITE},
    {'text': '每格 = A/B/C/D 等級', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'L4+ 最高 B 級（硬上限）', 'size': 10, 'color': ORANGE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '等級標準：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'A: WR≥60% + PF≥1.5 ✅', 'size': 10, 'color': GREEN},
    {'text': 'B: WR≥50% + PF≥1.2 ⚠️', 'size': 10, 'color': YELLOW},
    {'text': 'C: WR≥40% + PF≥1.0 🔶', 'size': 10, 'color': ORANGE},
    {'text': 'D: 其他 ❌', 'size': 10, 'color': RED_DANGER},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 歷史表現≠未來保證', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• L4+ B 級上限保守但合理', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# ============================================================
# SLIDE 10: Tab 13-16 — 波幅/馬丁風險/TPSL/Copy模擬
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📊 頁面功能 13-16：風險與模擬工具", "波幅分析 → 馬丁風險 → TP/SL → Copy 模擬")

# Tab 13: 波幅分析
add_rounded_rect(slide, 0.4, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 0.55, 1.35, 2.8, 0.35, "📏 波幅分析 (Volatility)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 0.55, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '層數波幅（CSV Max Pips）', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'ATR 交叉分析（本地 D1 數據）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '三組：低/中/高波幅（P33/P67）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Pip 乘數：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '5位小數對（AUDCAD）×10000', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'JPY 對 ×100', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'XAU/XAG ×0.1', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• ATR 來自 D1，短線精度低', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• pip 乘數配置可能有遺漏', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 14: 馬丁風險
add_rounded_rect(slide, 3.6, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 3.75, 1.35, 2.8, 0.35, "🎲 馬丁風險 (Martin Risk)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 3.75, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '評估馬丁格爾策略爆倉風險', 'size': 10, 'color': TEXT_WHITE},
    {'text': '風險等級 + 回歸率分析', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '風險等級：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'LOW: 最大層數 ≤ 3', 'size': 10, 'color': GREEN},
    {'text': 'MEDIUM: 最大層數 4-5', 'size': 10, 'color': YELLOW},
    {'text': 'HIGH: 最大層數 6-7', 'size': 10, 'color': ORANGE},
    {'text': 'CRITICAL: 最大層數 ≥ 8', 'size': 10, 'color': RED_DANGER},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '馬丁偵測類型：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'Classic: profit>0 且 pips<0', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Reverse: pips>0 且 profit<0', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Cost Killed: gross>0 且 net<0', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 51/69 信號有馬丁特徵！', 'size': 12, 'color': RED_DANGER, 'bold': True},
])

# Tab 15: TP/SL
add_rounded_rect(slide, 6.8, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 6.95, 1.35, 2.8, 0.35, "📐 TP/SL 分析", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 6.95, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '優化止損止盈建議', 'size': 10, 'color': TEXT_WHITE},
    {'text': '按貨幣對×層數 計算', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '公式（P85 基準）：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'TP = P85 of Max Pips', 'size': 10, 'color': TEXT_WHITE},
    {'text': '  → 85% 盈利交易曾到達此位', 'size': 10, 'color': TEXT_GRAY},
    {'text': 'SL = P85 of |Max Loss Pips|', 'size': 10, 'color': TEXT_WHITE},
    {'text': '  → 85% 交易回撤不超過此值', 'size': 10, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '小樣本回退：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'n < 30 → 用 global percentiles', 'size': 10, 'color': ORANGE},
    {'text': '69 個信號綜合數據', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 1863/2314 配對用 global fallback', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 無 Tick Data，觸發時間為估算', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 16: Copy 模擬
add_rounded_rect(slide, 10.0, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 10.15, 1.35, 2.8, 0.35, "💰 Copy 模擬 (Simulation)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 10.15, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '模擬「如果我咁 Copy 會點」', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '兩種模式：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'Copy on Profit（跟勝）', 'size': 10, 'color': GREEN},
    {'text': '  → 只跟 L1-L2，設獨立止損', 'size': 10, 'color': TEXT_WHITE},
    {'text': '  → 保守型，犧牲利潤換準確度', 'size': 10, 'color': TEXT_GRAY},
    {'text': 'Copy on Lose（跟虧）', 'size': 10, 'color': ORANGE},
    {'text': '  → 跟所有層數但設上限', 'size': 10, 'color': TEXT_WHITE},
    {'text': '  → 博反彈，控制最大虧損', 'size': 10, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '止損自動填入：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '從 TP/SL 分析建議值帶入', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 模擬結果≠實際交易結果', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 滑點、流動性未計入', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# ============================================================
# SLIDE 11: Tab 17-20 — 倉位明細/馬丁剖析V3/SET參數/歷史記錄
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📊 頁面功能 17-20：明細與系統工具", "倉位明細 → 馬丁剖析V3 → SET 參數 → 歷史記錄")

# Tab 17: 倉位明細
add_rounded_rect(slide, 0.4, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 0.55, 1.35, 2.8, 0.35, "📋 倉位明細 (Positions)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 0.55, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '所有交易逐筆列表', 'size': 10, 'color': TEXT_WHITE},
    {'text': '可排序、可篩選', 'size': 10, 'color': TEXT_WHITE},
    {'text': '點擊展開詳細資訊', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '顯示欄位：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'Open/Close Time, Symbol', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Type, Lots, Layer', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Pips, Profit, Max Pips', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Max Loss, Score', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 大量交易時渲染慢', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 無匯出功能（需手動 copy）', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 18: 馬丁剖析V3
add_rounded_rect(slide, 3.6, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 3.75, 1.35, 2.8, 0.35, "🔬 馬丁剖析V3 (Autopsy)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 3.75, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '6 部分完整馬丁剖析報告', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'Part 1: CCY×Direction 總覽', 'size': 10, 'color': TEXT_WHITE},
    {'text': '  EV$, Odds$, MFE/MAE', 'size': 10, 'color': TEXT_GRAY},
    {'text': 'Part 2: MFE/MAE 散點圖', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Part 3: TP/SL 混合方案', 'size': 10, 'color': TEXT_WHITE},
    {'text': '  Soft SL = MAE×1.2', 'size': 10, 'color': TEXT_GRAY},
    {'text': '  Hard SL = MaxMAE×1.3', 'size': 10, 'color': TEXT_GRAY},
    {'text': 'Part 4: 排行榜 (Rating+EV$)', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Part 5: 黑名單 (Danger Score)', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'Part 6: 恢復力分析', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 計算量大，大 CSV 較慢', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 19: SET 參數
add_rounded_rect(slide, 6.8, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 6.95, 1.35, 2.8, 0.35, "⚙️ SET 參數 (SET Params)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 6.95, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '解析 .set 設定檔並分類展示', 'size': 10, 'color': TEXT_WHITE},
    {'text': '參數分類：加倉/風控/TP/SL', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '解析：parseSET()', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': 'Key=Value 格式解析', 'size': 10, 'color': TEXT_WHITE},
    {'text': '自動識別 EA 家族', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'renderSETParams() 分類渲染', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• 新 EA 參數可能未涵蓋', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• UseAISignal=1 邏輯未知', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# Tab 20: 歷史記錄
add_rounded_rect(slide, 10.0, 1.3, 3.1, 5.9, BG_CARD)
add_text_box(slide, 10.15, 1.35, 2.8, 0.35, "📁 歷史記錄 (Archive)", font_size=13, color=ACCENT_RED, bold=True)
add_multi_text(slide, 10.15, 1.75, 2.8, 5.0, [
    {'text': '功能：', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '追蹤同一信號歷史表現變化', 'size': 10, 'color': TEXT_WHITE},
    {'text': 'CSV hash 偵測重複', 'size': 10, 'color': TEXT_WHITE},
    {'text': '趨勢迷你圖（勝率/PF/盈虧）', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '存儲：localStorage', 'size': 11, 'color': YELLOW, 'bold': True},
    {'text': '5MB 上限', 'size': 10, 'color': ORANGE},
    {'text': '智能標籤從檔名提取', 'size': 10, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '⚠️ 風險：', 'size': 11, 'color': RED_DANGER, 'bold': True},
    {'text': '• localStorage 單用戶限制', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 5MB 可能不夠大量分析', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 清瀏覽器數據會遺失', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
    {'text': '• 不支持跨設備同步', 'size': 10, 'color': RGBColor(0xFF, 0xCC, 0xCC)},
])

# ============================================================
# SLIDE 12: Backend Reports - Signal Ranking & Detailed Report
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📊 後端生成報告", "Signal Ranking 排名總表 + Detailed Comparison 詳細報告")

# Signal Ranking
add_rounded_rect(slide, 0.4, 1.3, 6.2, 5.9, BG_CARD)
add_text_box(slide, 0.6, 1.4, 5.8, 0.4, "🏆 Signal Ranking 排名總表", font_size=16, color=ACCENT_RED, bold=True)
add_multi_text(slide, 0.6, 1.9, 5.8, 5.0, [
    {'text': '生成方式：', 'size': 12, 'color': YELLOW, 'bold': True},
    {'text': 'python3 generate_signal_ranking.py', 'size': 11, 'color': GREEN},
    {'text': '→ 從 69 份 detailed reports 提取分數', 'size': 11, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '欄位說明（14 欄）：', 'size': 12, 'color': YELLOW, 'bold': True},
    {'text': 'Signal — 信號 ID（可點擊跳轉 AlgoForest）', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'Avg Score — DDE v3 平均分（越高越好）', 'size': 11, 'color': TEXT_WHITE},
    {'text': '⭐⭐⭐⭐ — 評分 ≥ 80 的組合數量', 'size': 11, 'color': TEXT_WHITE},
    {'text': '⭐⭐⭐⭐% — 高評分佔比', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'Trades / Win% / PF — 基本統計', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'Total Profit — 總盈利 ($)', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'TF — 時間框架（M30/H1/H4/D1+）', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'Cmp — 有效評分維度數', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'EA — EA 家族標籤（顏色編碼）', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'LV — 馬丁層數', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'Eq Max DD — 最大權益回撤', 'size': 11, 'color': TEXT_WHITE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '顏色分級：DD 🟢<$3K 🟡$3-6K 🔴>$6K', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'Score: 🟢≥90 🟢≥85 🟡≥75 🔴<75', 'size': 11, 'color': TEXT_WHITE},
])

# Detailed Report
add_rounded_rect(slide, 6.8, 1.3, 6.0, 5.9, BG_CARD)
add_text_box(slide, 7.0, 1.4, 5.6, 0.4, "📋 Detailed Comparison 詳細報告", font_size=16, color=ACCENT_RED, bold=True)
add_multi_text(slide, 7.0, 1.9, 5.6, 5.0, [
    {'text': '生成方式：', 'size': 12, 'color': YELLOW, 'bold': True},
    {'text': 'python3 generate_all_levels_from_csv.py --signal {ID}', 'size': 11, 'color': GREEN},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '報告結構（3 大區塊）：', 'size': 12, 'color': YELLOW, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '1️⃣ 分析摘要表', 'size': 12, 'color': ACCENT_RED, 'bold': True},
    {'text': '   貨幣對 × L1-L4+ 交易數 + 勝率', 'size': 11, 'color': TEXT_WHITE},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '2️⃣ 每個貨幣對（8-30 個）：', 'size': 12, 'color': ACCENT_RED, 'bold': True},
    {'text': '   📋 分析摘要表（CCY × Level）', 'size': 11, 'color': TEXT_WHITE},
    {'text': '   🎯 Copy Trade 建議引擎', 'size': 11, 'color': TEXT_WHITE},
    {'text': '      → 決策邏輯 + 信心度', 'size': 11, 'color': TEXT_GRAY},
    {'text': '   📈 值博率分析', 'size': 11, 'color': TEXT_WHITE},
    {'text': '      → EV + Kelly + Safety Margin', 'size': 11, 'color': TEXT_GRAY},
    {'text': '   🎰 馬丁層級深度分析', 'size': 11, 'color': TEXT_WHITE},
    {'text': '   CoP 評分表 (Wait 5/10/15/20)', 'size': 11, 'color': TEXT_WHITE},
    {'text': '   CoL 評分表 (Wait 10/15/20/25)', 'size': 11, 'color': TEXT_WHITE},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '3️⃣ Martin Detection', 'size': 12, 'color': ACCENT_RED, 'bold': True},
    {'text': '   Classic / Reverse / Cost Killed', 'size': 11, 'color': TEXT_WHITE},
])

# ============================================================
# SLIDE 13: System Risks Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "⚠️ 系統整體風險分析", "數據層面、計算層面、架構層面、決策層面")

risk_categories = [
    ("📥 數據風險", [
        "無 Tick Data — 無法精確知道入場後的實時價格走勢",
        "D1 精度不足 — 對 M5/H1 EA 市況分類不夠精細",
        "AutoLot 不確定 — 6 個信號的 lot 層級為推算",
        "CSV 時差 — Scraper vs 手動下載可能有幾行差異",
        "UseAISignal=1 — SMA SET 的 AI 信號邏輯未知",
    ]),
    ("🧮 計算風險", [
        "CoP 勝率永遠 100% — 只看盈利交易，20% 權重白送",
        "P85 基準假設穩定 — 市場 regime change 後失效",
        "小樣本回退 — 1863/2314 配對用 global fallback",
        "百分位外推 — 未來數據可能超出歷史範圍",
        "ETE 已廢棄但 DDE 也有盲區（dd_ratio capped 2.0）",
    ]),
    ("🏗️ 架構風險", [
        "單文件 HTML — index.html 已 246KB/5000+ 行，難維護",
        "localStorage 5MB 上限 — 大量分析可能爆限",
        "無多用戶 — 不支持跨設備同步",
        "localhost:8787 — 貨幣穿透功能依賴本地 API",
        "Cloudflare 攔截 — AlgoForest scraper 可能被擋",
    ]),
    ("💰 決策風險", [
        "歷史表現 ≠ 未來保證 — 所有評分基於回顧數據",
        "馬丁依賴 — 74% 信號有馬丁特徵，獲利可能是假象",
        "L4+ 爆倉高危 — 即使回歸率 70%，30% 失敗 = 巨虧",
        "模擬 ≠ 實際 — 滑點、流動性、延遲未計入",
        "黑天鵝事件 — 歷史數據無法預測極端市況",
    ]),
]

for i, (title, risks) in enumerate(risk_categories):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.3
    y = 1.3 + row * 3.1
    add_risk_card(slide, x, y, 6.1, 2.9, title, risks)

# ============================================================
# SLIDE 14: Complete Usage Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "🔄 完整使用流程 — 從數據到決策", "以「分析一個新 Signal」為例")

# Step-by-step flow
flow_data = [
    ("Step 1\n📥 取得數據", "AlgoForest 網站\n下載 CSV + SET\n或用 Scraper 自動", ACCENT_BLUE),
    ("Step 2\n🖥️ 上傳分析", "打開 index.html\n上傳 CSV（必須）\n上傳 SET（可選）", RGBColor(0x1B, 0x5E, 0x20)),
    ("Step 3\n📈 初步觀察", "收益曲線 Tab\n睇總體盈虧趨勢\n確認數據正常", ORANGE),
    ("Step 4\n🎯 質量評估", "質量評分 Tab\nEntry + Strategy Score\n識別弱項", ACCENT_RED),
    ("Step 5\n📚 深入分析", "層數/馬丁/TP/SL\n識別爆倉風險\n優化止損止盈", RGBColor(0x4A, 0x14, 0x8C)),
    ("Step 6\n🏆 Copy 決策", "Copy 推薦矩陣\nCopy 模擬驗證\n最終決定跟唔跟", RGBColor(0x00, 0x60, 0x64)),
]

for i, (title, desc, color) in enumerate(flow_data):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.5 + row * 3.0
    
    add_arrow_box(slide, x, y, 1.6, 0.7, title, color, font_size=10)
    add_rounded_rect(slide, x, y + 0.8, 3.9, 1.8, BG_CARD)
    add_text_box(slide, x + 0.15, y + 0.9, 3.6, 1.6, desc, font_size=12, color=TEXT_WHITE)

# Command reference
add_rounded_rect(slide, 0.5, 5.8, 12.3, 1.3, RGBColor(0x1B, 0x2A, 0x1B))
add_multi_text(slide, 0.7, 5.9, 11.9, 1.1, [
    {'text': '💻 快速指令參考：', 'size': 13, 'color': GREEN, 'bold': True},
    {'text': '後端分析：python3 generate_all_levels_from_csv.py --signal {ID} --csv samples/{ID}.csv', 'size': 11, 'color': TEXT_WHITE},
    {'text': '排名生成：python3 generate_signal_ranking.py                    |  Skill 觸發：跟丁蟹講「分析 signal {ID}」', 'size': 11, 'color': TEXT_WHITE},
])

# ============================================================
# SLIDE 15: TP/SL Calculation Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📐 TP/SL 計算原理深入", "P85 百分位法 + 小樣本回退 + 實例演算")

add_rounded_rect(slide, 0.5, 1.4, 6.0, 3.0, BG_CARD)
add_multi_text(slide, 0.7, 1.5, 5.6, 2.8, [
    {'text': '📐 TP/SL 建議公式', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'TP（Take Profit 建議值）', 'size': 13, 'color': GREEN, 'bold': True},
    {'text': 'TP = P85 of Max Pips（盈利交易）', 'size': 12, 'color': TEXT_WHITE},
    {'text': '含義：85% 嘅盈利交易曾到達此位置', 'size': 11, 'color': TEXT_GRAY},
    {'text': '即「85% 可達成」', 'size': 11, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': 'SL（Stop Loss 建議值）', 'size': 13, 'color': RED_DANGER, 'bold': True},
    {'text': 'SL = P85 of |Max Loss Pips|（所有交易）', 'size': 12, 'color': TEXT_WHITE},
    {'text': '含義：85% 嘅交易最大回撤不超過此值', 'size': 11, 'color': TEXT_GRAY},
    {'text': '即「85% 扛得住」', 'size': 11, 'color': TEXT_GRAY},
])

add_rounded_rect(slide, 6.8, 1.4, 6.0, 3.0, BG_CARD)
add_multi_text(slide, 7.0, 1.5, 5.6, 2.8, [
    {'text': '🔄 小樣本回退機制', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '問題：某 CCY×Level 可能只有幾筆交易', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'P85 在小樣本時不穩定', 'size': 12, 'color': ORANGE},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '解決方案：', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': 'n ≥ 30 → 用信號自身 percentiles', 'size': 12, 'color': GREEN},
    {'text': 'n < 30 → 自動混和 global percentiles', 'size': 12, 'color': ORANGE},
    {'text': 'Global = 69 個信號綜合數據', 'size': 11, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '統計：P85 足夠 193 / Hybrid 258 / Global 1863', 'size': 12, 'color': TEXT_WHITE},
    {'text': '⚠️ 大部分依賴 global fallback', 'size': 12, 'color': RED_DANGER},
])

# Example calculation
add_rounded_rect(slide, 0.5, 4.6, 12.3, 2.6, RGBColor(0x1B, 0x2A, 0x1B))
add_multi_text(slide, 0.7, 4.7, 11.9, 2.4, [
    {'text': '📝 實例演算：Signal 22200 / AUDCAD / L1 / CoP', 'size': 14, 'color': GREEN, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '盈利交易 Max Pips 排序（從小到大）：12, 18, 22, 25, 28, 30, 33, 35, 38, 40, 42, 45, 48, 50, 52, 55, 58, 60, 65, 70, 75, 80, 88, 95, 110, 125, 140, 165, 190, 220, 280, 350', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'n = 32 筆 ≥ 30，使用信號自身 percentiles', 'size': 11, 'color': GREEN},
    {'text': 'P85 位置 = 32 × 0.85 = 27.2 → 取第 28 個值 = 165 pips', 'size': 11, 'color': TEXT_WHITE},
    {'text': '→ TP 建議 = 165 pips（85% 盈利交易曾到達此位）', 'size': 12, 'color': ACCENT_RED, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '所有交易 |Max Loss Pips| 排序：...  P85 = 95 pips', 'size': 11, 'color': TEXT_WHITE},
    {'text': '→ SL 建議 = 95 pips（85% 交易回撤不超過此值）', 'size': 12, 'color': ACCENT_RED, 'bold': True},
    {'text': '→ R:R = 165/95 = 1.74  |  盈虧比合理 ✅', 'size': 12, 'color': GREEN, 'bold': True},
])

# ============================================================
# SLIDE 16: Martin Detection Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "🎲 馬丁偵測與風險分析", "74% 信號有馬丁特徵 — 這是最大的結構性風險")

# Martin types
for i, (mtype, cond, meaning, color, example) in enumerate([
    ("Classic Martin\n（經典馬丁）", 
     "profit > 0 且 pips < 0", 
     "方向錯了但靠加倉拉平成本獲利\n盈利是假象，依賴越來越大的注碼",
     ORANGE,
     "AUDCAD SELL: pips=-45 但 profit=$12\n→ 因為 L3 層 lots 大，少少反彈就賺"),
    ("Reverse Martin\n（反向馬丁）", 
     "pips > 0 且 profit < 0", 
     "方向對了但 swap/commission 吃掉利潤\n策略成本太高",
     YELLOW,
     "EURCHF BUY: pips=+15 但 profit=-$3\n→ 持倉 5 天，swap 吃掉所有利潤"),
    ("Cost Killed\n（成本殺手）", 
     "gross_profit > 0 且 net_profit < 0", 
     "毛利正但成本（commission+swap）吃掉淨利\n手續費/隔夜息太高",
     RED_DANGER,
     "XAUUSD BUY: gross=$8 但 net=-$2\n→ commission=$5 + swap=$5 吃掉"),
]):
    x = 0.4 + i * 4.2
    add_rounded_rect(slide, x, 1.4, 4.0, 3.3, BG_CARD)
    add_text_box(slide, x + 0.15, 1.45, 3.7, 0.45, mtype, font_size=13, color=color, bold=True)
    add_multi_text(slide, x + 0.15, 1.95, 3.7, 2.6, [
        {'text': f'條件：{cond}', 'size': 10, 'color': TEXT_GRAY},
        {'text': '', 'size': 4, 'color': TEXT_WHITE},
        {'text': meaning, 'size': 10, 'color': TEXT_WHITE},
        {'text': '', 'size': 6, 'color': TEXT_WHITE},
        {'text': f'例子：{example}', 'size': 9, 'color': TEXT_GRAY},
    ])

# Martin risk stats
add_rounded_rect(slide, 0.4, 4.9, 6.0, 2.3, RGBColor(0x2A, 0x15, 0x15))
add_multi_text(slide, 0.6, 5.0, 5.6, 2.1, [
    {'text': '⚠️ 馬丁風險統計（69 Signals）', 'size': 14, 'color': RED_DANGER, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '• 51/69 個信號（74%）有馬丁特徵', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 112 個 Classic Martin 貨幣對', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• L4+ 交易佔比越高 = 爆倉風險越大', 'size': 12, 'color': ORANGE},
    {'text': '• L4+ 回歸率 ≥ 70% 先算「可接受」', 'size': 12, 'color': TEXT_WHITE},
    {'text': '• 但 30% 失敗 = 巨額虧損', 'size': 12, 'color': RED_DANGER},
])

add_rounded_rect(slide, 6.8, 4.9, 6.0, 2.3, BG_CARD)
add_multi_text(slide, 7.0, 5.0, 5.6, 2.1, [
    {'text': '🛡️ 風險緩解措施', 'size': 14, 'color': GREEN, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '1. 優先 Copy L1-L2 交易（淺層安全）', 'size': 12, 'color': TEXT_WHITE},
    {'text': '2. 設硬止損（P85 of Max Loss Pips）', 'size': 12, 'color': TEXT_WHITE},
    {'text': '3. 用 Copy on Profit 策略（確認方向後跟）', 'size': 12, 'color': TEXT_WHITE},
    {'text': '4. 限制單一信號投入資金比例', 'size': 12, 'color': TEXT_WHITE},
    {'text': '5. 定期重新分析（市場環境變化）', 'size': 12, 'color': TEXT_WHITE},
])

# ============================================================
# SLIDE 17: Copy Trade Decision Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "🏆 Copy Trade 決策流程", "從分析到最終跟單建議的完整決策樹")

# Decision tree as flow
decisions = [
    (0.5, 1.5, 3.0, 0.7, "開始：上傳 CSV + SET", ACCENT_BLUE),
    (4.0, 1.5, 3.5, 0.7, "計算 DDE v3 Score", RGBColor(0x1B, 0x5E, 0x20)),
    (8.2, 1.5, 4.6, 0.7, "Avg Score ≥ 80？\n⭐⭐⭐⭐ 信號", ORANGE),
    
    (0.5, 2.8, 3.5, 0.7, "檢查馬丁依賴度", ACCENT_RED),
    (4.5, 2.8, 3.5, 0.7, "依賴 < 30% AND 勝率 > 60%?", RGBColor(0x1B, 0x5E, 0x20)),
    (8.5, 2.8, 4.3, 0.7, "✅ Copy on Profit\nWait = 最佳 CoP", GREEN),
    
    (0.5, 4.1, 3.5, 0.7, "依賴 ≥ 30% 有 CoL 數據？", ORANGE),
    (4.5, 4.1, 3.5, 0.7, "⚠️ Copy on Lose\nWait = 最佳 CoL", YELLOW),
    (8.5, 4.1, 4.3, 0.7, "❌ 不建議 Copy\n期望值 < 0.1R", RED_DANGER),
]

for x, y, w, h, text, color in decisions:
    add_arrow_box(slide, x, y, w, h, text, color, font_size=10)

# Arrows (text-based)
add_text_box(slide, 3.5, 1.7, 0.5, 0.3, "→", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 7.7, 1.7, 0.5, 0.3, "→", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2.2, 2.2, 0.3, 0.5, "↓", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 2.2, 0.3, 0.5, "↓ Yes    ↓ No", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2.2, 3.5, 0.3, 0.5, "↓", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 3.5, 0.3, 0.5, "↓ Yes    ↓ No", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Confidence levels
add_rounded_rect(slide, 0.5, 5.3, 12.3, 1.9, BG_CARD)
add_multi_text(slide, 0.7, 5.4, 11.9, 1.7, [
    {'text': '信心度分級：', 'size': 14, 'color': YELLOW, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '🟢 高信心 — 期望值 > 0.5R + 馬丁依賴 < 20% + 勝率 > 80%  →  可以放心 Copy，但永遠設止損', 'size': 12, 'color': GREEN},
    {'text': '🟡 中信心 — 介於高與低之間  →  可以 Copy 但要嚴格控制倉位，定期檢討', 'size': 12, 'color': YELLOW},
    {'text': '🔴 低信心 — 期望值 < 0.1R OR 馬丁依賴 > 70%  →  不建議 Copy，或只用極小倉位試水', 'size': 12, 'color': RED_DANGER},
])

# ============================================================
# SLIDE 18: Glossary & Quick Reference
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_header(slide, "📖 術語對照表與快速參考")

add_rounded_rect(slide, 0.4, 1.3, 6.2, 5.9, BG_CARD)
add_multi_text(slide, 0.6, 1.4, 5.8, 5.7, [
    {'text': '🔤 核心術語', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': 'DDE — Drawdown Efficiency 回撤效率', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'CoP — Copy on Profit 跟單盈利（確認方向後進場）', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'CoL — Copy on Lose 跟單虧損（博反彈進場）', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'PF — Profit Factor 盈利因子', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'WR — Win Rate 勝率', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'LV — Layer 層數（馬丁加倉層級）', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'MFE — Maximum Favorable Excursion 最大有利偏移', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'MAE — Maximum Adverse Excursion 最大不利偏移', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'EV — Expected Value 期望值', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'ATR — Average True Range 平均真實波幅', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'P85 — 第 85 百分位數', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'SET — EA 策略設定檔（key=value 格式）', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'DD — Drawdown 回撤', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'R:R — Reward:Risk Ratio 盈虧比', 'size': 12, 'color': TEXT_WHITE},
    {'text': 'Kelly — Kelly Criterion 凱利公式（最優注碼比例）', 'size': 12, 'color': TEXT_WHITE},
])

add_rounded_rect(slide, 6.8, 1.3, 6.0, 5.9, BG_CARD)
add_multi_text(slide, 7.0, 1.4, 5.6, 5.7, [
    {'text': '⚡ 快速操作指南', 'size': 16, 'color': ACCENT_RED, 'bold': True},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '前端（瀏覽器）：', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': 'URL: https://alvin-forex.github.io/trade-strategy-analyzer/', 'size': 11, 'color': GREEN},
    {'text': '或本地打開 docs/index.html', 'size': 11, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '後端分析（終端）：', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '# 單一信號分析', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'python3 generate_all_levels_from_csv.py \\', 'size': 10, 'color': GREEN},
    {'text': '  --signal {ID} --csv samples/{ID}.csv', 'size': 10, 'color': GREEN},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '# 生成排名總表', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'python3 generate_signal_ranking.py', 'size': 10, 'color': GREEN},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '# 批量分析所有信號', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'python3 batch_detailed_all.py', 'size': 10, 'color': GREEN},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '# 馬丁剖析 V3', 'size': 11, 'color': TEXT_WHITE},
    {'text': 'python3 generate_martin_autopsy_v3.py {csv}', 'size': 10, 'color': GREEN},
    {'text': '', 'size': 4, 'color': TEXT_WHITE},
    {'text': '# Skill 快捷指令', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '跟丁蟹講：「分析 signal {ID}」', 'size': 11, 'color': ACCENT_RED},
    {'text': '→ 自動下載 CSV + 分析 + Telegram 發送', 'size': 11, 'color': TEXT_GRAY},
    {'text': '', 'size': 6, 'color': TEXT_WHITE},
    {'text': '報告查詢命令：', 'size': 13, 'color': YELLOW, 'bold': True},
    {'text': '/history — 最近 5 條分析', 'size': 11, 'color': TEXT_WHITE},
    {'text': '/summary {id} — 查看詳情', 'size': 11, 'color': TEXT_WHITE},
    {'text': '/compare {signal} {v1} {v2} — 版本對比', 'size': 11, 'color': TEXT_WHITE},
])

# ============================================================
# SLIDE 19: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.8))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, 1, 1.5, 11, 0.5, "TRADE STRATEGY ANALYZER", font_size=16, color=TEXT_GRAY, bold=True)
add_text_box(slide, 1, 2.7, 11, 1.0, "系統說明書完", font_size=44, color=TEXT_WHITE, bold=True)
add_text_box(slide, 1, 3.8, 11, 0.6, "所有分析基於歷史數據，不構成投資建議。交易有風險，入市需謹慎。", font_size=16, color=YELLOW)

add_multi_text(slide, 1, 5.8, 11, 1.2, [
    {'text': f'📅 版本 v0.7  |  日期 {datetime.date.today().strftime("%Y-%m-%d")}  |  丁蟹 🦀 + Alvin', 'size': 14, 'color': TEXT_GRAY},
    {'text': '前端：https://alvin-forex.github.io/trade-strategy-analyzer/', 'size': 13, 'color': RGBColor(0x90, 0xCA, 0xF9)},
    {'text': '後端：localhost:8787  |  Skill：「分析 signal {ID}」', 'size': 13, 'color': TEXT_GRAY},
])

# ============================================================
# Save
# ============================================================
output_path = '/home/alvin/.openclaw/workspace/trade_strategy_analyzer/TSA_System_Manual.pptx'
prs.save(output_path)
print(f"✅ PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
